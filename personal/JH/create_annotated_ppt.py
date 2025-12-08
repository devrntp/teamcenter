# -*- coding: utf-8 -*-
"""
BMIDE 이미지에 주석 추가 및 PPT 생성 스크립트
"""

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 경로 설정
BASE_DIR = r"C:\Users\tuf\.gemini\antigravity\brain\22fb8f7e-e033-498c-9450-83782ea40d86"
OUTPUT_DIR = r"d:\dev\intellijProject\teamcenter\personal\JH"

# 이미지 경로
IMG1 = os.path.join(BASE_DIR, "uploaded_image_0_1765183907778.png")
IMG2 = os.path.join(BASE_DIR, "uploaded_image_1_1765183907778.png")
IMG3 = os.path.join(BASE_DIR, "uploaded_image_2_1765183907778.png")

# 색상 정의
RED = (255, 50, 50)
BLUE = (50, 100, 255)
GREEN = (50, 200, 100)
ORANGE = (255, 150, 50)
PURPLE = (150, 50, 200)
YELLOW = (255, 220, 50)
TEAL = (0, 155, 165)

def draw_rectangle(draw, coords, color, width=3):
    """사각형 그리기"""
    x1, y1, x2, y2 = coords
    for i in range(width):
        draw.rectangle([x1-i, y1-i, x2+i, y2+i], outline=color)

def draw_arrow(draw, start, end, color, width=3):
    """화살표 그리기"""
    draw.line([start, end], fill=color, width=width)
    # 화살표 머리
    import math
    dx = end[0] - start[0]
    dy = end[1] - start[1]
    angle = math.atan2(dy, dx)
    arrow_size = 15
    angle1 = angle + math.pi * 0.8
    angle2 = angle - math.pi * 0.8
    p1 = (end[0] + arrow_size * math.cos(angle1), end[1] + arrow_size * math.sin(angle1))
    p2 = (end[0] + arrow_size * math.cos(angle2), end[1] + arrow_size * math.sin(angle2))
    draw.polygon([end, p1, p2], fill=color)

def draw_numbered_circle(draw, center, number, color, radius=18):
    """번호가 있는 원 그리기"""
    x, y = center
    draw.ellipse([x-radius, y-radius, x+radius, y+radius], fill=color, outline=(255,255,255), width=2)
    # 번호 텍스트
    try:
        font = ImageFont.truetype("arial.ttf", 20)
    except:
        font = ImageFont.load_default()
    text = str(number)
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text((x - tw//2, y - th//2 - 2), text, fill=(255, 255, 255), font=font)

def draw_label(draw, pos, text, color, bg_color=(255, 255, 255, 200)):
    """라벨 텍스트 그리기"""
    try:
        font = ImageFont.truetype("malgun.ttf", 16)  # 한글 폰트
    except:
        try:
            font = ImageFont.truetype("arial.ttf", 14)
        except:
            font = ImageFont.load_default()
    
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    x, y = pos
    padding = 5
    
    # 배경 사각형
    draw.rectangle([x-padding, y-padding, x+tw+padding, y+th+padding], fill=(50, 50, 50, 230))
    draw.text((x, y), text, fill=color, font=font)

# =============================================================================
# 이미지 1: 전체 레이아웃 주석
# =============================================================================
def annotate_image1():
    img = Image.open(IMG1).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    # 1. 메뉴 바 영역
    draw_rectangle(draw, (0, 0, 1020, 35), RED, 3)
    draw_numbered_circle(draw, (35, 50), 1, RED)
    
    # 2. 툴바 영역
    draw_rectangle(draw, (0, 35, 1020, 75), BLUE, 3)
    draw_numbered_circle(draw, (35, 90), 2, BLUE)
    
    # 3. 좌측 트리 네비게이터
    draw_rectangle(draw, (0, 75, 245, 630), GREEN, 3)
    draw_numbered_circle(draw, (25, 350), 3, GREEN)
    
    # 4. 중앙 편집 영역
    draw_rectangle(draw, (250, 75, 780, 470), ORANGE, 3)
    draw_numbered_circle(draw, (515, 100), 4, ORANGE)
    
    # 5. 탭 메뉴
    draw_rectangle(draw, (270, 130, 730, 160), PURPLE, 3)
    draw_numbered_circle(draw, (750, 145), 5, PURPLE)
    
    # 6. Business Object Constants 테이블
    draw_rectangle(draw, (250, 470, 780, 630), TEAL, 3)
    draw_numbered_circle(draw, (515, 550), 6, TEAL)
    
    # 7. 하단 콘솔 영역
    draw_rectangle(draw, (0, 400, 245, 630), YELLOW, 3)
    draw_numbered_circle(draw, (123, 420), 7, YELLOW)
    
    output_path = os.path.join(OUTPUT_DIR, "bmide_annotated_1.png")
    img.save(output_path)
    return output_path

# =============================================================================
# 이미지 2: Business Object 상세 주석
# =============================================================================
def annotate_image2():
    img = Image.open(IMG2).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    # 1. Project
    draw_rectangle(draw, (430, 180, 790, 200), RED, 2)
    draw_numbered_circle(draw, (810, 190), 1, RED, 14)
    
    # 2. Name
    draw_rectangle(draw, (430, 195, 790, 215), BLUE, 2)
    draw_numbered_circle(draw, (810, 205), 2, BLUE, 14)
    
    # 3. Display Name
    draw_rectangle(draw, (430, 212, 790, 232), GREEN, 2)
    draw_numbered_circle(draw, (810, 222), 3, GREEN, 14)
    
    # 4. Storage Class
    draw_rectangle(draw, (430, 228, 550, 248), ORANGE, 2)
    draw_numbered_circle(draw, (570, 238), 4, ORANGE, 14)
    
    # 5. Parent-Item Revision-Form
    draw_rectangle(draw, (430, 245, 550, 295), PURPLE, 2)
    draw_numbered_circle(draw, (570, 270), 5, PURPLE, 14)
    
    # 6. Type 옵션
    draw_rectangle(draw, (350, 300, 650, 380), TEAL, 2)
    draw_numbered_circle(draw, (670, 340), 6, TEAL, 14)
    
    # 7. Teamcenter Component
    draw_rectangle(draw, (430, 375, 830, 395), YELLOW, 2)
    draw_numbered_circle(draw, (850, 385), 7, YELLOW, 14)
    
    output_path = os.path.join(OUTPUT_DIR, "bmide_annotated_2.png")
    img.save(output_path)
    return output_path

# =============================================================================
# 이미지 3: BMIDE 메뉴 주석
# =============================================================================
def annotate_image3():
    img = Image.open(IMG3).convert("RGBA")
    draw = ImageDraw.Draw(img)
    
    # 메뉴 항목들
    menu_items = [
        ((60, 95, 220, 115), 1, RED),      # New Model Element
        ((60, 115, 250, 135), 2, BLUE),    # New Composite Software Project
        ((60, 135, 120, 155), 3, GREEN),   # Find
        ((60, 165, 180, 185), 4, ORANGE),  # Save Data Model
        ((60, 185, 180, 205), 5, PURPLE),  # Reload Data Model
        ((60, 205, 220, 225), 6, TEAL),    # Generate Software Package
        ((60, 225, 180, 245), 7, YELLOW),  # Deploy Template
        ((60, 265, 140, 285), 8, RED),     # Live Update
        ((60, 285, 120, 305), 9, BLUE),    # Editors
    ]
    
    for coords, num, color in menu_items:
        draw_rectangle(draw, coords, color, 2)
        draw_numbered_circle(draw, (coords[2] + 20, (coords[1] + coords[3])//2), num, color, 12)
    
    output_path = os.path.join(OUTPUT_DIR, "bmide_annotated_3.png")
    img.save(output_path)
    return output_path

# =============================================================================
# PPT 생성
# =============================================================================
def create_annotated_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    DARK_BLUE = RGBColor(0, 60, 113)
    WHITE = RGBColor(255, 255, 255)
    SIEMENS_TEAL = RGBColor(0, 155, 165)
    BLACK = RGBColor(0, 0, 0)
    
    def add_slide_with_image(title, image_path, descriptions):
        """이미지와 설명이 있는 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 헤더
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_BLUE
        header.line.fill.background()
        
        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(10), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 이미지
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.2), Inches(1), width=Inches(7.5))
        
        # 설명 박스
        desc_left = Inches(7.9)
        desc_top = Inches(1)
        desc_box = slide.shapes.add_textbox(desc_left, desc_top, Inches(5.2), Inches(6.3))
        tf = desc_box.text_frame
        tf.word_wrap = True
        
        colors_map = {
            1: RGBColor(255, 50, 50),
            2: RGBColor(50, 100, 255),
            3: RGBColor(50, 200, 100),
            4: RGBColor(255, 150, 50),
            5: RGBColor(150, 50, 200),
            6: RGBColor(0, 155, 165),
            7: RGBColor(255, 180, 50),
            8: RGBColor(255, 50, 50),
            9: RGBColor(50, 100, 255),
        }
        
        for i, (num, text) in enumerate(descriptions):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"❶❷❸❹❺❻❼❽❾"[num-1] + " " + text
            p.font.size = Pt(14)
            p.font.color.rgb = colors_map.get(num, BLACK)
            p.font.bold = True
            p.space_after = Pt(8)
        
        return slide
    
    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Teamcenter BMIDE"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "화면 구성 가이드 - 이미지 주석 버전"
    p.font.size = Pt(24)
    p.font.color.rgb = SIEMENS_TEAL
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 1: 전체 레이아웃
    img1_path = annotate_image1()
    add_slide_with_image(
        "1. BMIDE 전체 화면 레이아웃",
        img1_path,
        [
            (1, "메뉴 바: File, BMIDE, Navigate, Search, Project, Run, Window, Help"),
            (2, "툴바: 저장, 열기, 실행, Launch Configurations 등"),
            (3, "트리 네비게이터: Business Objects, Classes, Navigator 탭"),
            (4, "편집 영역: 선택된 Business Object 상세 정보"),
            (5, "탭 메뉴: Main, Properties, Operations, Display Rules 등"),
            (6, "Business Object Constants: 상수 테이블"),
            (7, "Extensions/Console: 확장 및 출력 패널"),
        ]
    )
    
    # 슬라이드 2: Business Object 상세
    img2_path = annotate_image2()
    add_slide_with_image(
        "2. Business Object 상세 필드",
        img2_path,
        [
            (1, "Project: 현재 프로젝트명 (a2custom)"),
            (2, "Name: 오브젝트 내부 이름 (A2_custItem)"),
            (3, "Display Name: 화면 표시 이름 (cust Item)"),
            (4, "Storage Class: 저장소 클래스 타입"),
            (5, "Parent/Item Revision/Form: 상속 관계 및 연결 폼"),
            (6, "Type 옵션: Persistent, Abstract, Exportable 등"),
            (7, "Teamcenter Component: TC 컴포넌트 연결"),
        ]
    )
    
    # 슬라이드 3: BMIDE 메뉴
    img3_path = annotate_image3()
    add_slide_with_image(
        "3. BMIDE 메뉴 기능",
        img3_path,
        [
            (1, "New Model Element: 새 모델 요소 생성 (Ctrl+N)"),
            (2, "New Composite Software Project: 복합 프로젝트 생성"),
            (3, "Find: 요소 검색"),
            (4, "Save Data Model: 데이터 모델 저장 (Ctrl+Shift+S)"),
            (5, "Reload Data Model: 서버에서 다시 불러오기"),
            (6, "Generate Software Package: 배포 패키지 생성"),
            (7, "Deploy Template: 템플릿 서버 배포 (Ctrl+Shift+D)"),
            (8, "Live Update: 실시간 변경사항 반영"),
            (9, "Editors: 속성/LOV 에디터 등"),
        ]
    )
    
    # 저장
    output_path = os.path.join(OUTPUT_DIR, "Teamcenter_BMIDE_Guide_Annotated.pptx")
    prs.save(output_path)
    print(f"Annotated PPT saved to: {output_path}")
    return output_path

# 실행
if __name__ == "__main__":
    create_annotated_ppt()
