# -*- coding: utf-8 -*-
"""
Teamcenter BMIDE PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# RGBColor alias for convenience
RgbColor = RGBColor

# PPT 생성
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 색상 정의
SIEMENS_TEAL = RgbColor(0, 155, 165)  # Siemens 대표 색상
DARK_BLUE = RgbColor(0, 60, 113)
LIGHT_GRAY = RgbColor(240, 240, 240)
WHITE = RgbColor(255, 255, 255)
BLACK = RgbColor(0, 0, 0)

def add_title_slide(prs, title, subtitle=""):
    """표지 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경 사각형
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 서브타이틀
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(0, 200, 210)
        p.alignment = PP_ALIGN.CENTER
    
    # 하단 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(6), Inches(7.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SIEMENS_TEAL
    line.line.fill.background()
    
    return slide

def add_content_slide(prs, title, content_items, image_path=None):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 상단 헤더 바
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 내용 영역
    if image_path and os.path.exists(image_path):
        # 이미지가 있는 경우: 왼쪽 이미지, 오른쪽 텍스트
        slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.5), width=Inches(6.5))
        content_left = Inches(7)
        content_width = Inches(6)
    else:
        content_left = Inches(0.5)
        content_width = Inches(12.333)
    
    # 내용 텍스트
    content_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            p.font.color.rgb = item.get('color', BLACK)
            p.level = item.get('level', 0)
        else:
            p.text = "• " + item
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
        
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 상단 헤더 바
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 테이블 추가
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)).table
    
    # 헤더 스타일
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SIEMENS_TEAL
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 데이터 행
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    return slide

# =============================================================================
# 슬라이드 생성
# =============================================================================

# 1. 표지
add_title_slide(prs, 
    "Teamcenter BMIDE",
    "Business Modeler IDE 화면 구성 가이드")

# 2. 목차
add_content_slide(prs, "목차 (Contents)", [
    {"text": "1. BMIDE 개요", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "2. 전체 화면 레이아웃", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "3. 트리 네비게이터 (Tree Navigator)", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "4. Business Object 편집 영역", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "5. 탭 메뉴 상세", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "6. Business Object Constants", "size": 22, "bold": True, "color": DARK_BLUE},
    {"text": "7. BMIDE 메뉴 기능", "size": 22, "bold": True, "color": DARK_BLUE},
])

# 3. BMIDE 개요
add_content_slide(prs, "1. BMIDE 개요", [
    {"text": "BMIDE (Business Modeler IDE)란?", "size": 24, "bold": True, "color": SIEMENS_TEAL},
    "Teamcenter 데이터 모델을 정의하고 커스터마이징하는 통합 개발 환경",
    "Eclipse 기반의 개발 도구로, Teamcenter 서버의 비즈니스 오브젝트 구조를 설계",
    "",
    {"text": "주요 기능", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "Business Object (비즈니스 오브젝트) 생성 및 수정",
    "속성 (Property) 정의 및 관리",
    "LOV (List of Values) 설정",
    "Display Rules 및 Operation 정의",
    "데이터 모델 패키지 생성 및 배포",
])

# 4. 전체 화면 레이아웃
add_content_slide(prs, "2. 전체 화면 레이아웃", [
    {"text": "BMIDE 주요 영역 구성", "size": 24, "bold": True, "color": SIEMENS_TEAL},
    "",
    {"text": "① 상단 메뉴 바", "size": 18, "bold": True, "color": BLACK},
    "   File, BMIDE, Navigate, Search, Project, Run, Window, Help 메뉴 제공",
    "",
    {"text": "② 툴바 (Toolbar)", "size": 18, "bold": True, "color": BLACK},
    "   열기, 저장, 실행, Launch Configurations 등 빠른 접근 버튼",
    "",
    {"text": "③ 좌측 트리 네비게이터", "size": 18, "bold": True, "color": BLACK},
    "   Business Objects, Extensions, Classes 탭으로 구성된 탐색 패널",
    "",
    {"text": "④ 중앙 편집 영역", "size": 18, "bold": True, "color": BLACK},
    "   선택된 Business Object의 상세 정보를 편집하는 메인 작업 영역",
    "",
    {"text": "⑤ 하단 콘솔/출력 영역", "size": 18, "bold": True, "color": BLACK},
    "   Extensions, Outline, Console 탭으로 구성",
])

# 5. 트리 네비게이터 상세
add_content_slide(prs, "3. 트리 네비게이터 (Tree Navigator)", [
    {"text": "Business Objects 탭", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "프로젝트 내 모든 Business Object를 계층 구조로 표시",
    "폴더 구조: a2custom > Favorites > Business Objects",
    "",
    {"text": "주요 Business Object 유형", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "Item - 기본 아이템 오브젝트 (제품, 부품 등)",
    "Dataset - 첨부 파일 및 데이터셋 관리",
    "Folder - 폴더 구조 정의",
    "Form - 폼 오브젝트 정의",
    "Document - 문서 관리 오브젝트",
    "",
    {"text": "확장 패널", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "Classes 탭 - 클래스 계층 구조 확인",
    "Navigator 탭 - 빠른 탐색 기능",
])

# 6. 트리 내 오브젝트 종류
add_table_slide(prs, "3-1. 트리 내 주요 오브젝트 종류", 
    ["아이콘", "오브젝트 타입", "설명"],
    [
        ["📁", "Folder", "폴더 구조 정의, 오브젝트 그룹화"],
        ["📄", "Item", "기본 아이템 (부품, 제품, 문서 등)"],
        ["📋", "ItemRevision", "Item의 리비전 관리"],
        ["📎", "Dataset", "파일 첨부 및 데이터셋"],
        ["📝", "Form", "입력 폼 정의"],
        ["🏢", "Company", "회사/조직 정보"],
        ["👤", "Person", "사용자 정보"],
        ["🔧", "Tool", "도구 정의"],
        ["📐", "Design", "설계 정보"],
        ["🖼️", "Drawing", "도면 관리"],
    ])

# 7. Business Object 편집 영역
add_content_slide(prs, "4. Business Object 편집 영역", [
    {"text": "예시: A2_custItem 오브젝트", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "",
    {"text": "Details 섹션 필드 설명", "size": 20, "bold": True, "color": DARK_BLUE},
    "Project: a2custom (현재 프로젝트명)",
    "Name: A2_custItem (오브젝트 내부명)",
    "Display Name: cust Item (화면 표시명)",
    "Storage Class: A2_custItem (저장소 클래스)",
    "Parent: Item (부모 오브젝트 - 상속 관계)",
    "Item Revision: A2_custItemRevision (리비전 타입)",
    "Form: A2_custItemMaster (연결된 폼)",
    "Icon: Default (아이콘 설정)",
    "Type: Persistent (영구 저장 타입)",
])

# 8. Details 옵션 설명
add_table_slide(prs, "4-1. Details 섹션 옵션 설명",
    ["옵션", "설명", "체크 시 동작"],
    [
        ["Is Abstract?", "추상 클래스 여부", "인스턴스 생성 불가, 상속용으로만 사용"],
        ["Is Exportable?", "내보내기 가능 여부", "데이터 내보내기 시 포함"],
        ["Allow creating instances of Secondary Business Objects", "2차 오브젝트 생성 허용", "관련 오브젝트 자동 생성"],
        ["Store as lightweight object", "경량 오브젝트 저장", "성능 최적화를 위한 경량 저장"],
    ])

# 9. 탭 메뉴 개요
add_content_slide(prs, "5. 탭 메뉴 상세", [
    {"text": "Business Object 편집 탭 메뉴", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "",
    "Main - 기본 정보 및 속성 표시",
    "Properties - 오브젝트의 속성(Property) 정의",
    "Property Bulk Loaders - 대량 속성 로더 설정",
    "Operations - 비즈니스 로직 및 작업 정의",
    "Display Rules - 화면 표시 규칙 설정",
    "Deep Copy Rules - 깊은 복사 규칙 정의",
    "Alternate ID Rules - 대체 ID 규칙",
    "GRM Rules - 관계 매핑 규칙 (Generic Relationship Mapping)",
    "Operation Descriptor - 오퍼레이션 상세 정의",
])

# 10. 탭별 상세 기능
add_table_slide(prs, "5-1. 탭별 상세 기능",
    ["탭 이름", "주요 기능", "활용 예시"],
    [
        ["Main", "오브젝트 기본 정보 설정", "Name, Display Name, Parent 설정"],
        ["Properties", "속성 정의 및 관리", "사용자 정의 속성 추가, LOV 연결"],
        ["Property Bulk Loaders", "대량 데이터 로딩 설정", "Excel 데이터 일괄 등록"],
        ["Operations", "비즈니스 로직 정의", "생성/삭제/수정 시 동작 정의"],
        ["Display Rules", "UI 표시 규칙", "필드 표시/숨김, 읽기전용 설정"],
        ["Deep Copy Rules", "복사 시 하위 객체 포함 규칙", "Item 복사 시 Dataset 포함 여부"],
        ["GRM Rules", "관계 타입 매핑", "오브젝트 간 연결 관계 정의"],
    ])

# 11. Business Object Constants
add_content_slide(prs, "6. Business Object Constants", [
    {"text": "상수 테이블 구성", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "Business Object에 적용되는 상수값들을 정의하는 테이블",
    "",
    {"text": "테이블 컬럼 설명", "size": 20, "bold": True, "color": DARK_BLUE},
    "Name - 상수 이름 (예: Awb0AssociatedElem, Awb0EnableViewFor 등)",
    "Value - 설정값 (true/false 등)",
    "Overridden - 상위 클래스 값 오버라이드 여부",
    "Allow Modification - 수정 가능 여부",
    "Allow Override - 하위 클래스에서 오버라이드 가능 여부",
    "COTS - Commercial Off-The-Shelf (표준 제공 여부)",
    "Template - 적용 템플릿 (activeworks, aws2 등)",
])

# 12. BOC 주요 상수
add_table_slide(prs, "6-1. 주요 Business Object Constants",
    ["상수명", "설명", "일반적인 값"],
    [
        ["Awb0AssociatedElem", "AWC에서 연관 요소 표시", "true/false"],
        ["Awb0EnableViewForType", "타입별 뷰 활성화", "true/false"],
        ["Awb0BusinessObject", "AWC 비즈니스 오브젝트 설정", "aws2 템플릿"],
        ["Awb0DatasetTypeToBeCreated", "생성할 Dataset 타입", "타입명"],
        ["Awb0EnableCreateForType", "타입 생성 활성화", "true/false"],
        ["Awb0EnableSubmitFlow", "제출 워크플로우 활성화", "true/false"],
        ["Awb0SearchClassifySupport", "검색/분류 지원", "true/false"],
    ])

# 13. BMIDE 메뉴 기능
add_content_slide(prs, "7. BMIDE 메뉴 기능", [
    {"text": "BMIDE 메뉴 구성", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "",
    {"text": "주요 메뉴 항목", "size": 20, "bold": True, "color": DARK_BLUE},
    "New Model Element (Ctrl+N) - 새 모델 요소 생성",
    "New Composite Software Project - 복합 소프트웨어 프로젝트 생성",
    "Find - 요소 검색",
    "Save Data Model (Ctrl+Shift+S) - 데이터 모델 저장",
    "Reload Data Model - 데이터 모델 새로고침",
    "Generate Software Package - 소프트웨어 패키지 생성",
    "Deploy Template (Ctrl+Shift+D) - 템플릿 배포",
])

# 14. BMIDE 메뉴 상세
add_content_slide(prs, "7-1. BMIDE 메뉴 상세 기능", [
    {"text": "Live Update", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "실시간으로 서버에 변경사항 반영 (개발 중 테스트용)",
    "",
    {"text": "Editors", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "다양한 에디터 도구 접근 (Property Editor, LOV Editor 등)",
    "",
    {"text": "Organize Extensions", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "확장 모듈 관리 및 구성",
    "",
    {"text": "Upgrade Tools", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "버전 업그레이드 지원 도구",
    "",
    {"text": "Tools", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "추가 유틸리티 도구 모음",
    "",
    {"text": "Reports", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "데이터 모델 리포트 생성",
])

# 15. 버튼 및 툴바 설명
add_table_slide(prs, "8. 주요 버튼 및 툴바",
    ["버튼/아이콘", "기능", "단축키"],
    [
        ["💾 Save", "현재 편집 내용 저장", "Ctrl+S"],
        ["📁 Open", "파일/프로젝트 열기", "Ctrl+O"],
        ["🔄 Refresh", "뷰 새로고침", "F5"],
        ["▶ Run", "실행/배포", "-"],
        ["🔍 Search", "검색", "Ctrl+H"],
        ["📋 Copy", "복사", "Ctrl+C"],
        ["📝 Paste", "붙여넣기", "Ctrl+V"],
        ["↩ Undo", "실행 취소", "Ctrl+Z"],
        ["↪ Redo", "다시 실행", "Ctrl+Y"],
        ["🔧 Properties", "속성 창 열기", "Alt+Enter"],
    ])

# 16. 개발 워크플로우
add_content_slide(prs, "9. BMIDE 개발 워크플로우", [
    {"text": "일반적인 개발 순서", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "",
    {"text": "Step 1: 프로젝트 생성", "size": 18, "bold": True, "color": DARK_BLUE},
    "   New Composite Software Project로 커스텀 프로젝트 생성",
    "",
    {"text": "Step 2: Business Object 정의", "size": 18, "bold": True, "color": DARK_BLUE},
    "   Item 또는 다른 타입을 상속하여 새 오브젝트 생성",
    "",
    {"text": "Step 3: 속성 및 관계 설정", "size": 18, "bold": True, "color": DARK_BLUE},
    "   Properties 탭에서 사용자 정의 속성 추가",
    "",
    {"text": "Step 4: 데이터 모델 저장", "size": 18, "bold": True, "color": DARK_BLUE},
    "   Save Data Model로 변경사항 저장",
    "",
    {"text": "Step 5: 패키지 생성 및 배포", "size": 18, "bold": True, "color": DARK_BLUE},
    "   Generate Software Package → Deploy Template",
])

# 17. 마무리
add_content_slide(prs, "10. 요약 및 참고사항", [
    {"text": "BMIDE 핵심 포인트", "size": 22, "bold": True, "color": SIEMENS_TEAL},
    "Eclipse 기반으로 친숙한 IDE 환경 제공",
    "트리 구조로 Business Object 계층 관리",
    "다양한 탭을 통해 오브젝트의 모든 측면 설정 가능",
    "Live Update로 개발 중 실시간 테스트 가능",
    "",
    {"text": "주의사항", "size": 20, "bold": True, "color": RgbColor(200, 80, 80)},
    "데이터 모델 변경 시 반드시 저장 후 배포 필요",
    "운영 서버 반영 전 테스트 환경에서 충분한 검증 필요",
    "상속 관계 변경 시 하위 오브젝트에 영향 주의",
    "",
    {"text": "참고 문서", "size": 20, "bold": True, "color": SIEMENS_TEAL},
    "Teamcenter BMIDE Guide (공식 매뉴얼)",
    "Siemens Support Center (기술 지원)",
])

# PPT 저장
output_path = r"d:\dev\intellijProject\teamcenter\personal\JH\Teamcenter_BMIDE_Guide.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
